"""Build a PowerPoint presentation for the robust RAG project."""

from __future__ import annotations

from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
REPORTS = ROOT / "reports"
ASSET_DIR = Path(r"C:\Users\86159\.cursor\projects\d-360Downloads\assets")

FLOW_IMAGE = ASSET_DIR / (
    "c__Users_86159_AppData_Roaming_Cursor_User_workspaceStorage_"
    "3ff2ee083c5798b7e8e919b8bb352912_images_"
    "55d29c2c4b1bda1766f5a75733e54b1a-96877d73-6a0a-4863-b290-c93967d71861.png"
)
METHOD_IMAGE = ASSET_DIR / (
    "c__Users_86159_AppData_Roaming_Cursor_User_workspaceStorage_"
    "3ff2ee083c5798b7e8e919b8bb352912_images_"
    "3e1982d0aa7dab8153c66f5495ccfe18-5618ca58-634c-41f4-9b1b-148ac8e8330e.png"
)

BLUE = RGBColor(30, 72, 110)
LIGHT_BLUE = RGBColor(229, 239, 249)
GREEN = RGBColor(40, 130, 90)
RED = RGBColor(170, 55, 55)
ORANGE = RGBColor(190, 105, 35)
GRAY = RGBColor(90, 90, 90)
WHITE = RGBColor(255, 255, 255)


def set_run(run, size=20, bold=False, color=RGBColor(30, 30, 30), font="Microsoft YaHei"):
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_title(slide, title: str, subtitle: str | None = None):
    box = slide.shapes.add_textbox(Inches(0.35), Inches(0.18), Inches(12.65), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.alignment = PP_ALIGN.LEFT
    p.runs[0].font.name = "Microsoft YaHei"
    p.runs[0].font.size = Pt(22)
    p.runs[0].font.bold = True
    p.runs[0].font.color.rgb = BLUE
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.38), Inches(0.72), Inches(12.1), Inches(0.28))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        set_run(sp.runs[0], 10, False, GRAY)


def add_footer(slide, page: int):
    line = slide.shapes.add_shape(1, Inches(0.35), Inches(7.08), Inches(12.6), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(220, 225, 230)
    line.line.color.rgb = RGBColor(220, 225, 230)
    box = slide.shapes.add_textbox(Inches(0.35), Inches(7.1), Inches(12.6), Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = f"Robust RAG under Noisy Documents | {page}"
    p.alignment = PP_ALIGN.RIGHT
    set_run(p.runs[0], 8, False, GRAY)


def add_bullets(slide, items: list[str], x, y, w, h, size=18):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        set_run(p.runs[0], size, False)


def add_card(slide, x, y, w, h, title, body, color=LIGHT_BLUE):
    shape = slide.shapes.add_shape(5, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.color.rgb = RGBColor(210, 220, 230)
    box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.12), Inches(w - 0.36), Inches(h - 0.24))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    set_run(p.runs[0], 16, True, BLUE)
    p2 = tf.add_paragraph()
    p2.text = body
    p2.space_before = Pt(6)
    set_run(p2.runs[0], 12, False, RGBColor(45, 45, 45))


def add_table(slide, rows, x, y, w, h, font_size=10):
    table = slide.shapes.add_table(len(rows), len(rows[0]), Inches(x), Inches(y), Inches(w), Inches(h)).table
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            cell.margin_left = Inches(0.04)
            cell.margin_right = Inches(0.04)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    set_run(run, font_size, r == 0, WHITE if r == 0 else RGBColor(30, 30, 30))
            if r == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = BLUE
            elif r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(245, 248, 252)
    return table


def add_image_fit(slide, path: Path, x, y, w, h):
    with Image.open(path) as img:
        img_w, img_h = img.size
    box_ratio = w / h
    img_ratio = img_w / img_h
    if img_ratio >= box_ratio:
        final_w = w
        final_h = w / img_ratio
    else:
        final_h = h
        final_w = h * img_ratio
    final_x = x + (w - final_w) / 2
    final_y = y + (h - final_h) / 2
    slide.shapes.add_picture(str(path), Inches(final_x), Inches(final_y), width=Inches(final_w), height=Inches(final_h))


def add_kicker(slide, text: str, x=0.7, y=0.95, w=12.0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    p.text = text
    set_run(p.runs[0], 12, False, GRAY)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = RGBColor(248, 250, 252)
    slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.1)).fill.solid()
    slide.shapes[-1].fill.fore_color.rgb = BLUE
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(1.0))
    p = title.text_frame.paragraphs[0]
    p.text = "面向噪声文档的鲁棒 RAG 推理方法与评估"
    set_run(p.runs[0], 30, True, WHITE)
    sub = slide.shapes.add_textbox(Inches(0.75), Inches(1.65), Inches(11.8), Inches(0.6))
    p = sub.text_frame.paragraphs[0]
    p.text = "Evidence-Gated Iterative RAG · Controlled Noise · Failure Analysis"
    set_run(p.runs[0], 20, False, BLUE, "Arial")
    add_card(slide, 0.8, 2.7, 3.6, 1.4, "任务目标", "在 correct / noise / misinfo / contradictory 文档共存时，让 RAG 只基于可信证据作答。")
    add_card(slide, 4.85, 2.7, 3.6, 1.4, "核心方法", "EGI-RAG：证据门控、证据抽取、支持性验证；EGI-RAG+：冲突/误导风险控制。")
    add_card(slide, 8.9, 2.7, 3.6, 1.4, "汇报重点", "多维指标、噪声比例实验、错误案例分析、EGI-RAG+ 过度拒答问题。")
    add_footer(slide, 1)

    # 2
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "研究问题：RAG 不只是 Answer Accuracy")
    add_bullets(slide, [
        "真实检索结果中会混入普通噪声、误导答案和冲突文档。",
        "只看答案准确率，无法判断模型是否依赖了正确证据。",
        "因此评价扩展为：检索质量 + 证据忠实性 + 误导采纳 + 拒答能力。",
    ], 0.75, 1.25, 5.6, 3.1, 19)
    add_card(slide, 7.0, 1.25, 5.3, 0.95, "检索指标", "Recall@k, MRR, nDCG@k")
    add_card(slide, 7.0, 2.45, 5.3, 0.95, "证据指标", "Evidence F1, Strict Supported Rate")
    add_card(slide, 7.0, 3.65, 5.3, 0.95, "鲁棒性指标", "Misinfo Adoption Rate, Refusal F1")
    add_footer(slide, 2)

    # 3
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "实验全流程与成员分工", "数据构造、baseline、EGI-RAG、扩展评估与报告输出")
    add_image_fit(slide, FLOW_IMAGE, 0.35, 1.0, 12.65, 5.75)
    add_footer(slide, 3)

    # 4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "EGI-RAG 方法框架", "核心思想：先做证据过滤，再只基于可信证据生成与验证")
    add_image_fit(slide, METHOD_IMAGE, 0.4, 0.95, 12.5, 5.8)
    add_footer(slide, 4)

    # 5
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "方法实现：轻量工程版 EGI-RAG")
    add_card(slide, 0.65, 1.1, 3.7, 1.35, "1. Retrieval / Rerank", "top_k 到 top_n；本次补跑使用 lexical fallback，避免本地神经依赖冲突。")
    add_card(slide, 4.8, 1.1, 3.7, 1.35, "2. Evidence Gate", "LLM 判断 supportive / partial / irrelevant / misleading；EGI+ 增加 contradictory。")
    add_card(slide, 8.95, 1.1, 3.7, 1.35, "3. Evidence-only Generation", "只把 evidence spans 交给生成器，不直接拼接全文上下文。")
    add_card(slide, 2.1, 3.25, 4.2, 1.45, "4. Verification / Refusal", "答案必须 supported；unsupported/conflict/insufficient 时拒答。")
    add_card(slide, 7.1, 3.25, 4.2, 1.45, "实现边界", "图中的多轮 Iterative Corrector 和数值 doc score 是完整设计，主实验采用一轮轻量实现。", RGBColor(250, 240, 230))
    add_footer(slide, 5)

    # 6
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "全量主结果：EGI-RAG 与 EGI-RAG+")
    rows = [
        ["Dataset", "Method", "N", "AnsAcc", "Token F1", "Misinfo", "StrictSup"],
        ["RGB", "EGI-RAG", "300", "0.9200", "0.8127", "0.0000", "0.9667"],
        ["RGB", "EGI-RAG+", "300", "0.8600", "0.7616", "0.0000", "0.8767"],
        ["RAMDocs", "EGI-RAG", "500", "0.5320", "0.5381", "0.1100", "0.7120"],
        ["RAMDocs", "EGI-RAG+", "500", "0.2260", "0.2307", "0.0080", "0.2620"],
    ]
    add_table(slide, rows, 0.7, 1.2, 11.9, 2.2, 12)
    add_card(slide, 1.0, 4.05, 5.4, 1.25, "关键观察", "EGI-RAG 在 RGB 上准确率和支持性最好；RAMDocs 更强调误导采纳风险。")
    add_card(slide, 6.9, 4.05, 5.4, 1.25, "EGI-RAG+ 的定位", "Misinfo 从 0.1100 降到 0.0080，但 Accuracy 从 0.5320 降到 0.2260。它是风险控制型改进，不是 accuracy 改进。", RGBColor(250, 240, 230))
    add_footer(slide, 6)

    # 7
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Controlled Noise：RGB 代表性比例")
    rows = [
        ["Setting", "Method", "N", "AnsAcc", "StrictSup"],
        ["0% front", "EGI-RAG", "300", "0.9600", "0.9967"],
        ["0% front", "EGI-RAG+", "300", "0.9633", "0.9867"],
        ["60% front", "EGI-RAG", "300", "0.9300", "0.9667"],
        ["60% front", "EGI-RAG+", "300", "0.9233", "0.9533"],
        ["100% front", "EGI-RAG", "300", "0.0400", "0.1233"],
        ["100% front", "EGI-RAG+", "300", "0.0233", "0.0733"],
    ]
    add_table(slide, rows, 0.65, 1.05, 12.0, 3.2, 11)
    add_bullets(slide, [
        "0% 和 60% 普通噪声下，EGI 系列仍能保持较高 Answer Accuracy。",
        "100% 噪声时 Recall@5=0，模型没有正确证据，Accuracy 大幅下降是预期结果。",
        "全噪声压力测试更应该看拒答能力，而不是只看 Accuracy。",
    ], 1.0, 4.6, 11.5, 1.3, 15)
    add_footer(slide, 7)

    # 8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Controlled Noise：RAMDocs 代表性比例")
    rows = [
        ["Setting", "Method", "N", "AnsAcc", "Misinfo", "StrictSup"],
        ["0% front", "EGI-RAG", "365", "0.6055", "0.0438", "0.7260"],
        ["0% front", "EGI-RAG+", "290", "0.3655", "0.0103", "0.4276"],
        ["60% front", "EGI-RAG", "118", "0.5424", "0.1017", "0.7288"],
        ["60% front", "EGI-RAG+", "100", "0.3100", "0.0100", "0.3500"],
        ["100% front", "EGI-RAG", "100", "0.0100", "0.3200", "0.4100"],
        ["100% front", "EGI-RAG+", "100", "0.0000", "0.2200", "0.2900"],
    ]
    add_table(slide, rows, 0.55, 1.05, 12.2, 3.15, 10)
    add_card(slide, 0.9, 4.65, 5.6, 1.25, "结论 1", "RAMDocs 的主要风险不是检索不到，而是 high-overlap misinfo 被模型采纳。")
    add_card(slide, 6.85, 4.65, 5.6, 1.25, "结论 2", "EGI-RAG+ 显著压低 Misinfo，但强冲突门控造成过度拒答。", RGBColor(250, 240, 230))
    add_footer(slide, 8)

    # 9
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Custom Noise：逻辑缺失与高重叠误导")
    rows = [
        ["Method", "N", "AnsAcc", "Token F1", "Misinfo", "EvF1", "StrictSup"],
        ["Naive RAG", "60", "0.6500", "0.4694", "0.0333", "0.0000", "0.0000"],
        ["Rerank RAG", "60", "0.6500", "0.4703", "0.0167", "0.0000", "0.0000"],
        ["CRAG-lite", "60", "0.6333", "0.3511", "0.0167", "0.0000", "0.0000"],
        ["EGI-RAG", "60", "0.7833", "0.7755", "0.0167", "0.8611", "0.8667"],
    ]
    add_table(slide, rows, 0.55, 1.05, 12.2, 2.45, 11)
    add_bullets(slide, [
        "Custom Noise 专门构造 logic-gap、value-swap 和 high-overlap irrelevant 文档。",
        "在 Recall 已接近 1.0 的情况下，差异主要来自模型是否能过滤不完整/误导证据。",
        "EGI-RAG 通过 evidence spans 降低全文噪声干扰。",
    ], 0.9, 4.1, 11.6, 1.4, 15)
    add_footer(slide, 9)

    # 10
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "Failure Analysis：答错的问题为什么会错？")
    add_card(slide, 0.6, 1.0, 3.9, 1.15, "检索不到正确证据", "100% noise 时 R@5/MRR 为 0，生成器无可靠依据。", RGBColor(245, 248, 252))
    add_card(slide, 4.8, 1.0, 3.9, 1.15, "强误导文档共现", "RAMDocs 中 wrong answer 文档与问题高度相关。", RGBColor(245, 248, 252))
    add_card(slide, 9.0, 1.0, 3.7, 1.15, "证据抽取不完整", "抽到主体但遗漏时间、地点、数值等关键限定。", RGBColor(245, 248, 252))
    add_card(slide, 2.3, 3.0, 4.0, 1.2, "过度拒答", "EGI-RAG+ 对 misleading / contradictory 惩罚太硬。", RGBColor(250, 240, 230))
    add_card(slide, 7.1, 3.0, 4.0, 1.2, "字符串评估误差", "日期格式、别名、简称可能让语义正确答案被判错。", RGBColor(250, 240, 230))
    add_footer(slide, 10)

    # 11
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "EGI-RAG+ 为什么“变差”？")
    add_bullets(slide, [
        "它的目标是降低 misinformation adoption，不是直接提高 Answer Accuracy。",
        "当前策略是 hard gate：一旦检测到 misleading/contradictory 且 supportive evidence 不够强，就倾向拒答。",
        "这会减少错误答案采纳，但也会误拒答一批本来可回答的样本。",
        "因此报告中应表述为：风险控制有效，但覆盖率和准确率受损。",
    ], 0.9, 1.25, 11.7, 3.0, 18)
    add_card(slide, 1.3, 4.75, 4.9, 1.0, "有效的一面", "RAMDocs 全量 Misinfo: 0.1100 → 0.0080", RGBColor(235, 248, 240))
    add_card(slide, 7.1, 4.75, 4.9, 1.0, "代价", "RAMDocs 全量 AnsAcc: 0.5320 → 0.2260", RGBColor(252, 238, 238))
    add_footer(slide, 11)

    # 12
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, "结论与后续改进")
    add_bullets(slide, [
        "鲁棒 RAG 的关键不是读更多文档，而是只基于可信证据作答。",
        "EGI-RAG 在 RGB 和 Custom Noise 上表现稳定，证据链指标明显更好。",
        "RAMDocs 证明 high-overlap misinfo 是更危险的噪声类型。",
        "EGI-RAG+ 需要从 hard gate 改为 soft threshold：按 supportive 证据强度和冲突强度加权。",
        "后续应增加 LLM judge / 人工复核，缓解字符串评估误差。",
    ], 0.85, 1.15, 11.9, 3.6, 18)
    add_card(slide, 2.0, 5.25, 9.3, 0.85, "Takeaway", "Evidence first, answer second: trustworthy RAG should answer only from supported evidence.", RGBColor(235, 245, 252))
    add_footer(slide, 12)

    out = REPORTS / "rag_noise_robustness_presentation.pptx"
    prs.save(out)
    return out


def build() -> Path:
    """Build a cleaner, presentation-first deck."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def blank():
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(250, 252, 254)
        return slide

    # 1. Title
    slide = blank()
    band = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
    band.fill.solid()
    band.fill.fore_color.rgb = BLUE
    band.line.color.rgb = BLUE
    title = slide.shapes.add_textbox(Inches(0.75), Inches(0.36), Inches(11.8), Inches(0.65))
    p = title.text_frame.paragraphs[0]
    p.text = "面向噪声文档的鲁棒 RAG 推理"
    set_run(p.runs[0], 30, True, WHITE)
    subtitle = slide.shapes.add_textbox(Inches(0.8), Inches(1.55), Inches(11.8), Inches(0.55))
    p = subtitle.text_frame.paragraphs[0]
    p.text = "Evidence-Gated Iterative RAG · 检索维度评估 · 错误案例分析"
    set_run(p.runs[0], 20, False, BLUE)
    add_card(slide, 0.9, 2.75, 3.55, 1.25, "问题", "候选文档混入 noise / misinfo / contradictory 信息。")
    add_card(slide, 4.9, 2.75, 3.55, 1.25, "方法", "证据门控 + 证据抽取 + 支持性验证。")
    add_card(slide, 8.9, 2.75, 3.55, 1.25, "结论", "EGI-RAG 稳定；EGI-RAG+ 降低误导但过度拒答。")
    add_footer(slide, 1)

    # 2. Research problem
    slide = blank()
    add_title(slide, "研究问题：RAG 不只是 Answer Accuracy")
    add_kicker(slide, "真实检索结果可能包含正确证据、普通噪声、高重叠误导文档和冲突文档。")
    add_card(slide, 0.8, 1.65, 3.8, 1.25, "检索质量", "正确证据是否进入 top-k？\nRecall@k / MRR / nDCG")
    add_card(slide, 4.75, 1.65, 3.8, 1.25, "证据忠实性", "答案是否由证据支持？\nEvidence F1 / StrictSup")
    add_card(slide, 8.7, 1.65, 3.8, 1.25, "抗误导能力", "是否采纳 wrong answer？\nMisinfo / Refusal F1")
    add_bullets(
        slide,
        [
            "Naive/Rerank 可能检索到正确文档，但仍被 misinfo 带偏。",
            "100% noise 场景下，合理行为应从“作答”转向“拒答”。",
            "EGI-RAG 关注：先找到可信证据，再生成答案。",
        ],
        1.1,
        4.0,
        11.2,
        1.6,
        17,
    )
    add_footer(slide, 2)

    # 3. Flow image
    slide = blank()
    add_title(slide, "实验全流程与成员分工")
    add_kicker(slide, "从数据转换、controlled noise 构造，到 baseline / EGI-RAG 运行与扩展评估。")
    add_image_fit(slide, FLOW_IMAGE, 0.45, 1.05, 12.45, 5.75)
    add_footer(slide, 3)

    # 4. Method image
    slide = blank()
    add_title(slide, "EGI-RAG 方法框架")
    add_kicker(slide, "核心：检索重排后进行证据级判断，只基于 supportive evidence 生成并验证答案。")
    add_image_fit(slide, METHOD_IMAGE, 0.45, 1.05, 12.45, 5.75)
    add_footer(slide, 4)

    # 5. Full results
    slide = blank()
    add_title(slide, "全量主结果：EGI-RAG vs. EGI-RAG+")
    rows = [
        ["Dataset", "Method", "N", "AnsAcc", "Token F1", "Misinfo", "StrictSup"],
        ["RGB", "EGI-RAG", "300", "0.9200", "0.8127", "0.0000", "0.9667"],
        ["RGB", "EGI-RAG+", "300", "0.8600", "0.7616", "0.0000", "0.8767"],
        ["RAMDocs", "EGI-RAG", "500", "0.5320", "0.5381", "0.1100", "0.7120"],
        ["RAMDocs", "EGI-RAG+", "500", "0.2260", "0.2307", "0.0080", "0.2620"],
    ]
    add_table(slide, rows, 0.8, 1.25, 11.8, 2.35, 12)
    add_card(slide, 1.05, 4.15, 5.35, 1.25, "EGI-RAG", "RGB 上答案与证据支持性最高；RAMDocs 上降低误导但仍保持一定覆盖。")
    add_card(slide, 6.95, 4.15, 5.35, 1.25, "EGI-RAG+", "Misinfo 显著下降，但 hard gate 导致过度拒答，Accuracy 下降。", RGBColor(250, 240, 230))
    add_footer(slide, 5)

    # 6. Controlled trends
    slide = blank()
    add_title(slide, "噪声比例实验：代表性趋势")
    rows = [
        ["Setting", "Method", "N", "AnsAcc", "Misinfo", "StrictSup"],
        ["RGB 0%", "EGI-RAG", "300", "0.9600", "0.0000", "0.9967"],
        ["RGB 60%", "EGI-RAG", "300", "0.9300", "0.0000", "0.9667"],
        ["RGB 100%", "EGI-RAG", "300", "0.0400", "0.0000", "0.1233"],
        ["RAM 60%", "EGI-RAG", "118", "0.5424", "0.1017", "0.7288"],
        ["RAM 60%", "EGI-RAG+", "100", "0.3100", "0.0100", "0.3500"],
        ["RAM 100%", "EGI-RAG+", "100", "0.0000", "0.2200", "0.2900"],
    ]
    add_table(slide, rows, 0.75, 1.15, 11.9, 3.15, 11)
    add_bullets(
        slide,
        [
            "RGB：普通噪声下 EGI-RAG 仍稳定；100% noise 因无正确证据而大幅下降。",
            "RAMDocs：misinfo 是主要风险；EGI-RAG+ 降低误导采纳，但拒答过多。",
        ],
        0.95,
        4.75,
        11.6,
        1.0,
        15,
    )
    add_footer(slide, 6)

    # 7. Custom noise and failure
    slide = blank()
    add_title(slide, "Custom Noise 与错误来源")
    rows = [
        ["Method", "N", "AnsAcc", "Token F1", "Misinfo", "EvF1"],
        ["Naive", "60", "0.6500", "0.4694", "0.0333", "0.0000"],
        ["Rerank", "60", "0.6500", "0.4703", "0.0167", "0.0000"],
        ["CRAG", "60", "0.6333", "0.3511", "0.0167", "0.0000"],
        ["EGI-RAG", "60", "0.7833", "0.7755", "0.0167", "0.8611"],
    ]
    add_table(slide, rows, 0.75, 1.05, 11.8, 2.25, 12)
    add_card(slide, 0.95, 3.8, 3.55, 1.0, "检索失败", "100% noise 时没有正确证据。")
    add_card(slide, 4.85, 3.8, 3.55, 1.0, "误导共现", "正确证据与 wrong answer 文档同时相关。")
    add_card(slide, 8.75, 3.8, 3.55, 1.0, "过度拒答", "EGI-RAG+ 阈值过硬，覆盖率下降。", RGBColor(250, 240, 230))
    add_footer(slide, 7)

    # 8. Why EGI+ worse
    slide = blank()
    add_title(slide, "为什么 EGI-RAG+ 看起来更差？")
    add_kicker(slide, "它是风险控制型改进，不是直接优化 Answer Accuracy。")
    add_card(slide, 1.0, 1.65, 5.2, 1.35, "收益", "RAMDocs 全量 Misinfo Adoption\n0.1100 → 0.0080", RGBColor(235, 248, 240))
    add_card(slide, 7.05, 1.65, 5.2, 1.35, "代价", "RAMDocs 全量 Answer Accuracy\n0.5320 → 0.2260", RGBColor(252, 238, 238))
    add_bullets(
        slide,
        [
            "当前 EGI-RAG+ 使用 hard gate：遇到 misleading/contradictory 就更倾向拒答。",
            "这能减少 wrong answer 采纳，但也误拒答了一些有 supportive evidence 的样本。",
            "后续应改成 soft threshold：按证据强度和冲突强度加权。",
        ],
        1.05,
        4.05,
        11.4,
        1.5,
        17,
    )
    add_footer(slide, 8)

    # 9. Conclusion
    slide = blank()
    add_title(slide, "结论")
    add_bullets(
        slide,
        [
            "鲁棒 RAG 的关键不是读更多文档，而是只基于可信证据作答。",
            "EGI-RAG 在 RGB 与 Custom Noise 上提升明显，并提供可解释 evidence spans。",
            "RAMDocs 证明 high-overlap misinfo 比普通 noise 更危险。",
            "EGI-RAG+ 有效降低误导采纳，但需要缓解过度拒答。",
        ],
        0.95,
        1.35,
        11.8,
        3.1,
        20,
    )
    add_card(slide, 2.05, 5.25, 9.1, 0.85, "Takeaway", "Evidence first, answer second: trustworthy RAG should answer only from supported evidence.", RGBColor(235, 245, 252))
    add_footer(slide, 9)

    out = REPORTS / "rag_noise_robustness_presentation.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    path = build()
    print(f"[OK] wrote {path}")
